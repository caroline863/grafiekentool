"""
Excel -> PowerPoint Report Generator
=====================================
Streamlit app: upload crosstab Excel, configure via Regie Tabel,
generate styled PowerPoint with python-pptx.
"""
import io
import os
import re
import base64
from collections import OrderedDict
from pathlib import Path
import pandas as pd
import streamlit as st
import requests
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
)
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
# ── Gemini API setup ──
_gemini_api_key = os.environ.get("GEMINI_API_KEY")
# =====================================================================
# CONSTANTS
# =====================================================================
DEBUG_MODE = False  # Zet op True als je wil debuggen
FONT_NAME = "Avenir Next LT Pro"
FONT_SIZE = Pt(10)
# Patterns in question+answer text that indicate a row should be SKIPPED
# entirely (these rows act as block separators — they trigger _flush).
SKIP_PATTERNS = [
    "comparisons of column proportions",
    "results are based on",
]
# Question-level skip: if the QUESTION text itself equals one of these,
# the entire question block is discarded.
SKIP_QUESTION_EXACT = {"topbox", "bottombox"}
# Answer-level utility rows: these are collected into the block (so we
# can extract n=) but are DROPPED from the chart data afterwards.
# IMPORTANT: do NOT put these in SKIP_ANSWER_EXACT — that would flush
# the block prematurely and orphan any rows that come after.
DROP_ANSWER_OPTIONS = {"n", "%", "topbox", "bottombox", "top2box", "bottom2box"}
STACKED_KEYWORDS = [
    "eens", "oneens", "goed", "slecht", "vaak", "nooit",
    "belangrijk", "tevreden", "ontevreden",
]
STACKED_PREFIXES = ["zeer ", "heel ", "helemaal "]
STACKED_COLOR_MAP = {
    "helemaal eens":          "#39B54A",
    "helemaal mee eens":      "#39B54A",
    "zeer goed":              "#39B54A",
    "zeer tevreden":          "#39B54A",
    "mee eens":               "#A8CD66",
    "eens":                   "#A8CD66",
    "goed":                   "#A8CD66",
    "tevreden":               "#A8CD66",
    "neutraal":               "#FFC04D",
    "niet eens, niet oneens": "#FFC04D",
    "niet eens/niet oneens":  "#FFC04D",
    "mee oneens":             "#F28E2B",
    "oneens":                 "#F28E2B",
    "slecht":                 "#F28E2B",
    "ontevreden":             "#F28E2B",
    "helemaal oneens":        "#E1001A",
    "helemaal mee oneens":    "#E1001A",
    "zeer slecht":            "#E1001A",
    "zeer ontevreden":        "#E1001A",
    "weet ik niet":           "#808080",
    "weet niet":              "#808080",
    "geen mening":            "#808080",
}
STACKED_SCALE_COLORS = ["#39B54A", "#A8CD66", "#FFC04D", "#F28E2B", "#E1001A", "#808080"]
BAR_PRIMARY   = "#C60651"
BAR_SECONDARY = "#FF85A2"
BAR_GREY      = "#D3D3D3"
BOTTOM_LABELS = {
    "weet ik niet", "weet niet", "geen van bovenstaande", "geen mening",
}
PENULTIMATE_LABELS = {
    "anders, namelijk", "anders", "overig",
    "anders, namelijk...", "anders, namelijk:",
    "anders, namelijk …", "anders, namelijk…",
    "een ander netwerk, namelijk",
    "een ander netwerk, namelijk ...",
    "een ander netwerk, namelijk…",
}
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY  = RGBColor(0x80, 0x80, 0x80)

# 5-point scale patterns for stacked chart detection
FIVE_POINT_SCALE_SETS = [
    {"helemaal mee eens", "mee eens", "neutraal", "mee oneens", "helemaal mee oneens"},
    {"helemaal eens", "eens", "neutraal", "oneens", "helemaal oneens"},
    {"zeer tevreden", "tevreden", "neutraal", "ontevreden", "zeer ontevreden"},
    {"zeer goed", "goed", "neutraal", "slecht", "zeer slecht"},
    {"heel vaak", "vaak", "soms", "zelden", "nooit"},
]

# =====================================================================
# HELPERS
# =====================================================================
def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))
def _is_empty(val) -> bool:
    if val is None:
        return True
    if isinstance(val, float) and pd.isna(val):
        return True
    s = str(val).replace("\xa0", " ").strip().lower()
    return s in ("", "nan", "none")
def _clean_answer(raw: str) -> str:
    """Aggressively clean an answer string."""
    return (
        str(raw)
        .replace("\xa0", " ")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
        .strip()
    )
def _is_utility_row(answer: str) -> bool:
    """Check if this is a utility row (n, %, topbox, bottombox, etc.).
    Handles multi-line values like 'n\\n%', hidden characters, padding.
    """
    clean = _clean_answer(answer).lower()
    # Exact match
    if clean in DROP_ANSWER_OPTIONS:
        return True
    # Multi-line: split on newlines and check if ALL parts are utility
    parts = [p.strip() for p in clean.split("\n") if p.strip()]
    if parts and all(p in DROP_ANSWER_OPTIONS for p in parts):
        return True
    # Combined like "n %" or "n%"
    if re.match(r'^[n%\s]+$', clean) and len(clean) <= 5:
        return True
    return False
def _row_has_n(answer: str) -> bool:
    """Check if this answer row represents an 'n' (sample size) row."""
    clean = _clean_answer(answer).lower()
    if clean == "n":
        return True
    # Multi-line: one of the parts is "n"
    parts = [p.strip() for p in clean.split("\n") if p.strip()]
    return "n" in parts
def _should_skip_row(q_cell: str, a_cell: str) -> bool:
    """Check if a row should be SKIPPED entirely (acts as block separator).
    Only use this for truly separating content like footnotes.
    Do NOT use this for n/% /topbox/bottombox — those must stay in the
    block so n= can be extracted before they are filtered out.
    """
    combined = (q_cell + " " + a_cell).lower().strip()
    for pat in SKIP_PATTERNS:
        if pat in combined:
            return True
    return False
def _set_font(font, size=None, bold=False, color=None, name=FONT_NAME):
    font.name = name
    font.size = size if size else FONT_SIZE
    font.bold = bold
    if color:
        font.color.rgb = color


def _get_blank_layout(prs):
    """Get a blank slide layout safely. Tries layout index 6 (blank),
    then falls back to the last available layout."""
    layouts = prs.slide_layouts
    # Try common blank layout indices
    for idx in (6, 5, len(layouts) - 1, 0):
        if idx < len(layouts):
            return layouts[idx]
    return layouts[0]


def generate_significance_text(vraag, bevindingen_bullets):
    """
    Calls Gemini to rewrite rule-based significance bullets into
    professional Dutch report language with percentages included.
    Returns a string, or None if the call fails.
    """
    if not bevindingen_bullets:
        return None
    prompt = f"""Je bent een marktonderzoeker die bevindingen rapporteert in professionele stijl.
Vraag: "{vraag}"
Hieronder staan significante doelgroepverschillen in ruwe vorm:
{chr(10).join(bevindingen_bullets)}
Schrijf dit om naar 1 tot 3 vloeiende Nederlandse zinnen in rapportagestijl.
Regels:
- Vermeld altijd de percentages in de zin, bijvoorbeeld: "Mannen vinden vaker dan vrouwen dat X (45% versus 32%)"
- Groepeer per dimensie: eerst geslacht, dan leeftijd, dan opleiding
- Geen opsommingstekens, gewone lopende zinnen
- Geen inleiding of afsluiting, alleen de bevindingen zelf
- Maximaal 3 zinnen totaal
- Professioneel en feitelijk, geen interpretaties"""
    if not _gemini_api_key:
        return None
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={_gemini_api_key}"
        payload = {"contents": [{"parts": [{"text": prompt}]}]}
        response = requests.post(url, json=payload, timeout=15)
        result = response.json()
        return result["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        print(f"Gemini API error: {e}")
        return None
def _parse_sig_table(sig_rows_list, data_cols):
    """Parse a significance table from captured rows.
    Returns dict with dimension_groups and data_rows, or None."""
    if not sig_rows_list:
        return None
    df = pd.DataFrame(sig_rows_list)
    letter_pattern = re.compile(r'^\(([A-Z])\)$')
    # Find the letter row: contains cells like "(A)", "(B)", "(C)"
    letter_row_idx = None
    for idx in range(len(df)):
        matches = 0
        for col in data_cols:
            if col not in df.columns:
                continue
            val = str(df.iloc[idx][col]).strip()
            if letter_pattern.match(val):
                matches += 1
        if matches >= 2:
            letter_row_idx = idx
            break
    if letter_row_idx is None:
        return None
    # Detect dimension groups from the row above the letter row
    dimension_groups = []
    dim_row_idx = letter_row_idx - 1 if letter_row_idx > 0 else None
    if dim_row_idx is not None:
        current_dim = None
        current_cols = []
        for col in data_cols:
            if col not in df.columns:
                continue
            dim_val = str(df.iloc[dim_row_idx][col]).strip()
            if dim_val and dim_val.lower() not in ("nan", ""):
                if current_dim is not None:
                    dimension_groups.append({
                        "name": current_dim, "columns": current_cols, "letters": {}
                    })
                current_dim = dim_val
                current_cols = [col]
            else:
                if current_dim is not None:
                    current_cols.append(col)
        if current_dim is not None:
            dimension_groups.append({
                "name": current_dim, "columns": current_cols, "letters": {}
            })
    if not dimension_groups:
        dimension_groups = [{"name": "Totaal", "columns": data_cols, "letters": {}}]
    # Assign letters to columns within each dimension group
    for group in dimension_groups:
        for col in group["columns"]:
            if col not in df.columns:
                continue
            val = str(df.iloc[letter_row_idx][col]).strip()
            m = letter_pattern.match(val)
            if m:
                group["letters"][m.group(1)] = col
    # Parse data rows (after letter row)
    sig_data_rows = []
    for idx in range(letter_row_idx + 1, len(df)):
        answer = str(df.iloc[idx].get("_answer_", "")).strip()
        if not answer or answer.lower() in ("nan", ""):
            continue
        row_sigs = {}
        for col in data_cols:
            if col not in df.columns:
                continue
            cell = str(df.iloc[idx][col]).strip()
            if cell and cell.lower() not in ("nan", ""):
                letters = re.findall(r'[A-Z]', cell)
                if letters:
                    row_sigs[col] = letters
        if row_sigs:
            sig_data_rows.append({"answer": answer, "significances": row_sigs})
    if not sig_data_rows:
        return None
    return {"dimension_groups": dimension_groups, "data_rows": sig_data_rows}
def _generate_sig_bullets(sig_data, question_info):
    """Generate rule-based significance bullets from parsed significance data."""
    if not sig_data:
        return []
    df = question_info["df"]
    dimension_groups = sig_data["dimension_groups"]
    data_rows = sig_data["data_rows"]
    # Build reverse lookup: col_name -> (dimension_name, letter)
    col_to_info = {}
    dim_letter_to_col = {}
    for group in dimension_groups:
        dim_letter_to_col[group["name"]] = group["letters"]
        for letter, col_name in group["letters"].items():
            col_to_info[col_name] = (group["name"], letter)
    # Build answer -> col -> percentage lookup
    pct_lookup = {}
    for _, row in df.iterrows():
        ans = str(row["_answer_"]).strip()
        pct_lookup[ans] = {}
        for col in df.columns:
            if col.startswith("_"):
                continue
            try:
                pct_lookup[ans][col] = float(row[col]) if pd.notna(row[col]) else 0.0
            except (ValueError, TypeError):
                pass
    # Aggregate: (higher_col, lower_col) -> [(answer, h_pct, l_pct)]
    pair_findings = {}
    for row in data_rows:
        answer = row["answer"]
        for col_name, lower_letters in row["significances"].items():
            if col_name not in col_to_info:
                continue
            dim_name, _ = col_to_info[col_name]
            h_pct = pct_lookup.get(answer, {}).get(col_name, 0)
            for lower_letter in lower_letters:
                lower_col = dim_letter_to_col.get(dim_name, {}).get(lower_letter)
                if not lower_col:
                    continue
                l_pct = pct_lookup.get(answer, {}).get(lower_col, 0)
                key = (col_name, lower_col)
                if key not in pair_findings:
                    pair_findings[key] = []
                pair_findings[key].append((answer, h_pct, l_pct))
    # Generate bullets with percentages
    bullets = []
    for (higher_col, lower_col), findings in pair_findings.items():
        if len(findings) == 1:
            ans, h_pct, l_pct = findings[0]
            bullets.append(
                f"{higher_col} ({h_pct:.0f}%) scoort significant hoger dan "
                f"{lower_col} ({l_pct:.0f}%) op '{ans}'"
            )
        else:
            parts = []
            for ans, h_pct, l_pct in findings:
                parts.append(f"'{ans}' ({h_pct:.0f}% vs {l_pct:.0f}%)")
            bullets.append(
                f"{higher_col} scoort significant hoger dan {lower_col} op "
                + " en ".join(parts)
            )
    return bullets
# =====================================================================
# 1. PARSE EXCEL
# =====================================================================
def parse_excel(uploaded_file) -> tuple[OrderedDict, list[str]]:
    raw = pd.read_excel(uploaded_file, header=None, dtype=str)
    row0 = raw.iloc[0].fillna("").astype(str).str.strip()
    row1 = raw.iloc[1].fillna("").astype(str).str.strip()
    # Build column names from the two header rows
    col_names: list[str] = []
    current_main = ""
    # Characters that act as placeholder sub-headers (not real names)
    placeholder_subs = {"nan", "", "_", "-", "–", "—"}
    for idx, (main, sub) in enumerate(zip(row0, row1)):
        if main and main.lower() != "nan":
            current_main = main
        if sub and sub.lower() not in placeholder_subs:
            col_names.append(sub)
        elif current_main:
            col_names.append(current_main)
        else:
            col_names.append(f"_col{idx}_")
    col_names[0] = "_question_"
    col_names[1] = "_answer_"
    # De-duplicate data column names
    seen_names: dict[str, int] = {}
    for i, name in enumerate(col_names):
        if i < 2:
            continue
        if name in seen_names:
            seen_names[name] += 1
            col_names[i] = f"{name}_{seen_names[name]}"
        else:
            seen_names[name] = 0
    raw.columns = col_names
    data = raw.iloc[2:].reset_index(drop=True)
    data_cols = [c for c in col_names if c not in ("_question_", "_answer_")]
    questions: OrderedDict = OrderedDict()
    seen_questions: set = set()
    current_q: str | None = None
    current_rows: list[dict] = []
    last_flushed_q: str | None = None
    sig_capture: bool = False
    sig_rows_buf: list[dict] = []
    sig_question: str | None = None
    def _flush():
        nonlocal current_q, current_rows, last_flushed_q
        if current_q is None or not current_rows:
            current_q = None
            current_rows = []
            return
        q_key = current_q.strip()
        if q_key.lower().strip() in SKIP_QUESTION_EXACT:
            current_q = None
            current_rows = []
            return
        if q_key in seen_questions:
            current_q = None
            current_rows = []
            return
        seen_questions.add(q_key)
        df_block = pd.DataFrame(current_rows)
        # Force _answer_ to clean strings
        df_block["_answer_"] = (
            df_block["_answer_"]
            .fillna("")
            .astype(str)
            .apply(_clean_answer)
        )
        # ── Extract n= value BEFORE dropping utility rows ──
        n_value = "?"
        for idx_row in range(len(df_block)):
            ans = df_block.iloc[idx_row]["_answer_"]
            if _row_has_n(ans):
                # Try first available data column for the n value
                for dc in data_cols:
                    if dc not in df_block.columns:
                        continue
                    raw_n = df_block.iloc[idx_row][dc]
                    if _is_empty(raw_n):
                        continue
                    n_str = (
                        str(raw_n)
                        .replace("%", "")
                        .replace(",", "")
                        .replace("\xa0", "")
                        .strip()
                    )
                    try:
                        n_num = int(float(n_str))
                        if n_num > 0:
                            n_value = str(n_num)
                            break
                    except (ValueError, TypeError):
                        continue
                if n_value != "?":
                    break
        # ── Drop utility rows (n, %, topbox, bottombox) ──
        keep_mask = ~df_block["_answer_"].apply(_is_utility_row)
        df_clean = df_block[keep_mask].copy()
        # ── Convert data columns to numeric percentages ──
        for col in data_cols:
            if col not in df_clean.columns:
                continue
            # Clean string values
            series = (
                df_clean[col]
                .fillna("")
                .astype(str)
                .str.replace("%", "", regex=False)
                .str.replace(",", ".", regex=False)
                .str.replace("\xa0", "", regex=False)
                .str.strip()
            )
            df_clean[col] = pd.to_numeric(series, errors="coerce")
        # If all values are in 0-1 range (decimals), scale to 0-100
        all_vals = []
        for col in data_cols:
            if col in df_clean.columns:
                all_vals.extend(df_clean[col].dropna().tolist())
        if all_vals and max(all_vals) <= 1.0:
            for col in data_cols:
                if col in df_clean.columns:
                    df_clean[col] = df_clean[col] * 100
        answer_opts = df_clean["_answer_"].dropna().str.strip().tolist()
        # Extra filter: remove any answer that is empty after cleaning
        answer_opts = [a for a in answer_opts if a]
        questions[q_key] = {
            "df": df_clean.reset_index(drop=True),
            "n_value": n_value,
            "answer_options": answer_opts,
        }
        last_flushed_q = q_key
        current_q = None
        current_rows = []
    def _flush_sig():
        nonlocal sig_capture, sig_rows_buf, sig_question
        if sig_question and sig_rows_buf and sig_question in questions:
            parsed = _parse_sig_table(sig_rows_buf, data_cols)
            if parsed:
                bullets = _generate_sig_bullets(parsed, questions[sig_question])
                if bullets:
                    # Store raw bullets; Gemini rewrite happens at generation time
                    questions[sig_question]["significantie_bullets"] = bullets
                    questions[sig_question]["heeft_significantie"] = True
        sig_capture = False
        sig_rows_buf = []
        sig_question = None
    for _, row in data.iterrows():
        q_cell = str(row["_question_"]).strip() if not _is_empty(row["_question_"]) else ""
        a_cell = str(row["_answer_"]).strip() if not _is_empty(row["_answer_"]) else ""
        combined_lower = (q_cell + " " + a_cell).lower().strip()
        # Skip footnote rows (these ARE block separators)
        if _should_skip_row(q_cell, a_cell):
            _flush()
            current_q = None
            if sig_capture:
                _flush_sig()
            # Check if this is a significance table header
            if "comparisons of column proportions" in combined_lower:
                sig_capture = True
                sig_question = last_flushed_q
                sig_rows_buf = []
            continue
        # During significance capture mode
        if sig_capture:
            if not q_cell and not a_cell:
                _flush_sig()
                continue
            if q_cell:
                # New question starts — end sig capture, fall through
                _flush_sig()
            else:
                sig_rows_buf.append(row.to_dict())
                continue
        # Empty row = end of block
        if not q_cell and not a_cell:
            _flush()
            continue
        # New question starts
        if q_cell:
            _flush()
            current_q = q_cell
        # Add row to current block (including n/% — they'll be filtered later)
        if current_q is not None and a_cell:
            current_rows.append(row.to_dict())
    _flush()
    _flush_sig()
    # Set heeft_significantie=False for questions without significance data
    for q_key in questions:
        if "heeft_significantie" not in questions[q_key]:
            questions[q_key]["heeft_significantie"] = False
            questions[q_key]["significantie_bullets"] = []
    return questions, data_cols
# =====================================================================
# 2. AUTO-DETECT CHART TYPE
# =====================================================================
def detect_chart_type(answer_options: list[str]) -> str:
    """Default to bar chart. Only use stacked for clear 5-point scale questions."""
    answer_lower = {a.strip().lower() for a in answer_options if a.strip()}
    # Check against known 5-point scale sets
    for scale_set in FIVE_POINT_SCALE_SETS:
        # If answers contain at least 4 of 5 scale items, it's a scale question
        overlap = answer_lower & scale_set
        if len(overlap) >= 4:
            return "100% Gestapeld horizontaal"
    # Also check: if we have exactly 5-7 answers and they ALL match stacked color map
    if 4 <= len(answer_lower) <= 7:
        all_match = all(
            a in STACKED_COLOR_MAP or
            any(a in k or k in a for k in STACKED_COLOR_MAP)
            for a in answer_lower
            if a not in {"weet ik niet", "weet niet", "geen mening"}
        )
        non_grey = answer_lower - {"weet ik niet", "weet niet", "geen mening"}
        if all_match and len(non_grey) >= 4:
            return "100% Gestapeld horizontaal"
    return "Staafdiagram"
# =====================================================================
# 3. POWERPOINT GENERATION
# =====================================================================
def _sort_bar_df(df: pd.DataFrame, data_cols: list[str],
                 grey_labels: set[str] | None = None) -> pd.DataFrame:
    if not data_cols:
        return df
    first_col = data_cols[0]
    answers_lower = df["_answer_"].str.strip().str.lower()
    bottom_set = grey_labels if grey_labels is not None else BOTTOM_LABELS
    bottom_mask = answers_lower.isin(bottom_set)
    penult_mask = answers_lower.apply(
        lambda x: any(x.startswith(p) or x == p for p in PENULTIMATE_LABELS)
    )
    normal_mask = ~(bottom_mask | penult_mask)
    df_normal = df[normal_mask].sort_values(by=first_col, ascending=True)
    df_penult = df[penult_mask]
    df_bottom = df[bottom_mask]
    return pd.concat([df_bottom, df_penult, df_normal], ignore_index=True)
def _set_chart_title(chart, question: str, n_value: str, group_id: str):
    """Add question + basis as chart title (two paragraphs)."""
    chart.has_title = True
    tf = chart.chart_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    _set_font(p.font, size=Pt(10), bold=True, color=DARK_GREY)
    p.alignment = PP_ALIGN.CENTER
    basis_text = f"Basis: {group_id} (n={n_value})" if group_id else f"Basis: totaal (n={n_value})"
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = basis_text
    _set_font(run2.font, size=Pt(10), bold=False, color=MID_GREY)
    p2.alignment = PP_ALIGN.CENTER
def _clean_axes(chart):
    for axis_attr in ("value_axis", "category_axis"):
        ax = getattr(chart, axis_attr, None)
        if ax is None:
            continue
        ax.has_major_gridlines = False
        ax.has_minor_gridlines = False
        try:
            ax.major_tick_mark = 2  # NONE
            ax.minor_tick_mark = 2
        except Exception:
            pass
        try:
            ax.format.line.fill.background()
        except Exception:
            pass
    if chart.value_axis:
        chart.value_axis.visible = False
def _add_significance_textbox(slide, significantie_tekst):
    """Add significance text box below the chart on a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(0.39), Inches(6.72), Inches(12.55), Inches(0.65)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run1 = para.add_run()
    run1.text = "Significante verschillen: "
    run1.font.bold = True
    run1.font.size = Pt(9)
    run1.font.color.rgb = RGBColor(0x50, 0x50, 0x50)
    run2 = para.add_run()
    run2.text = significantie_tekst
    run2.font.bold = False
    run2.font.size = Pt(9)
    run2.font.color.rgb = RGBColor(0x50, 0x50, 0x50)
def _build_bar_slide(prs, question: str, info: dict,
                     selected_cols: list[str], group_id: str,
                     grey_label: str = "", significantie_tekst: str | None = None):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    df = info["df"].copy()
    n_value = info["n_value"]
    keep = ["_answer_"] + [c for c in selected_cols if c in df.columns]
    df = df[[c for c in keep if c in df.columns]]
    chart_data_cols = [c for c in selected_cols if c in df.columns]
    if df.empty or not chart_data_cols:
        return
    # Extra safety: drop any remaining utility rows
    df = df[~df["_answer_"].apply(_is_utility_row)].copy()
    if df.empty:
        return
    grey_set = {grey_label.strip().lower()} if grey_label.strip() else BOTTOM_LABELS
    df = _sort_bar_df(df, chart_data_cols, grey_labels=grey_set)
    categories = df["_answer_"].tolist()
    chart_data = CategoryChartData()
    chart_data.categories = categories
    # Reverse series order so visual top-to-bottom matches selection order
    for col in reversed(chart_data_cols):
        chart_data.add_series(col, df[col].fillna(0).tolist())
    # Position chart within slide guides (~0.5" margins)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(6.7),
        chart_data,
    )
    chart = chart_frame.chart
    _set_chart_title(chart, question, n_value, group_id)
    plot = chart.plots[0]
    plot.gap_width = 60
    _clean_axes(chart)
    if chart.category_axis:
        cat_font = chart.category_axis.tick_labels.font
        _set_font(cat_font, size=Pt(10), color=DARK_GREY)
    answers_lower = df["_answer_"].str.strip().str.lower().tolist()
    colors = [BAR_PRIMARY, BAR_SECONDARY, "#FF6B81", "#A855F7"]
    num_series = len(chart_data_cols)
    for s_idx, series in enumerate(plot.series):
        # Map colors so first selected column gets primary color
        color_idx = num_series - 1 - s_idx
        base_color = colors[color_idx % len(colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = hex_to_rgb(base_color)
        for pt_idx, ans in enumerate(answers_lower):
            if ans in grey_set:
                pt = series.points[pt_idx]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = hex_to_rgb(BAR_GREY)
        # Data labels — percentage at end of each bar
        series.has_data_labels = True
        dl = series.data_labels
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False
        dl.font.name = FONT_NAME
        dl.font.size = Pt(10)
        dl.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        dl.number_format = '0"%"'
        dl.number_format_is_linked = False
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
    if num_series <= 1:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        _set_font(chart.legend.font, size=Pt(10))
    if significantie_tekst:
        _add_significance_textbox(slide, significantie_tekst)
def _build_stacked_slide(prs, question: str, info: dict,
                         selected_cols: list[str], group_id: str,
                         significantie_tekst: str | None = None):
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    df = info["df"].copy()
    n_value = info["n_value"]
    keep = ["_answer_"] + [c for c in selected_cols if c in df.columns]
    df = df[[c for c in keep if c in df.columns]]
    chart_data_cols = [c for c in selected_cols if c in df.columns]
    if df.empty or not chart_data_cols:
        return
    # Extra safety: drop any remaining utility rows
    df = df[~df["_answer_"].apply(_is_utility_row)].copy()
    if df.empty:
        return
    chart_data = CategoryChartData()
    # Reverse category order so visual top-to-bottom matches selection order
    display_cols = list(reversed(chart_data_cols))
    chart_data.categories = display_cols
    answer_labels = df["_answer_"].str.strip().tolist()
    for _, row in df.iterrows():
        label = str(row["_answer_"]).strip()
        vals = [float(row[c]) if pd.notna(row[c]) else 0.0 for c in display_cols]
        chart_data.add_series(label, vals)
    # Position chart within slide guides (~0.5" margins)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100,
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(6.7),
        chart_data,
    )
    chart = chart_frame.chart
    _set_chart_title(chart, question, n_value, group_id)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.overlap = 100
    _clean_axes(chart)
    if chart.category_axis:
        cat_font = chart.category_axis.tick_labels.font
        _set_font(cat_font, size=Pt(10), color=DARK_GREY)
    num_answers = len(answer_labels)
    for s_idx, series in enumerate(plot.series):
        label = answer_labels[s_idx] if s_idx < num_answers else ""
        matched = _match_stacked_color(label.lower().strip())
        color = matched if matched else _get_stacked_position_color(s_idx, num_answers)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = hex_to_rgb(color)
        series.has_data_labels = True
        dl = series.data_labels
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False
        dl.font.name = FONT_NAME
        dl.font.size = Pt(10)
        dl.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        dl.number_format = '0"%"'
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.CENTER
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _set_font(chart.legend.font, size=Pt(10))
    if significantie_tekst:
        _add_significance_textbox(slide, significantie_tekst)
def _match_stacked_color(text: str) -> str | None:
    text = text.strip().lower()
    if text in STACKED_COLOR_MAP:
        return STACKED_COLOR_MAP[text]
    for key, color in STACKED_COLOR_MAP.items():
        if key in text or text in key:
            return color
    return None
def _get_stacked_position_color(idx: int, total: int) -> str:
    """Fallback scale color based on position (green -> red) for stacked charts."""
    if total <= 1:
        return STACKED_SCALE_COLORS[0]
    scale = STACKED_SCALE_COLORS[:5]  # exclude grey
    pos = idx / (total - 1)
    scale_idx = min(int(pos * (len(scale) - 1) + 0.5), len(scale) - 1)
    return scale[scale_idx]
def _build_grouped_stacked_slide(prs, group_questions: list[tuple[str, dict]],
                                  selected_cols: list[str], group_id: str,
                                  significantie_tekst: str | None = None):
    """Grouped stacked 100% chart: each question becomes a category row,
    answer options become the stacked segments.  Uses first selected column."""
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    first_col = selected_cols[0] if selected_cols else None
    if not first_col:
        return
    # Collect all unique answer options (order from first question)
    all_answers: list[str] = []
    seen_answers: set[str] = set()
    for _, info in group_questions:
        df = info["df"].copy()
        df = df[~df["_answer_"].apply(_is_utility_row)].copy()
        for ans in df["_answer_"].str.strip().tolist():
            if ans and ans not in seen_answers:
                all_answers.append(ans)
                seen_answers.add(ans)
    if not all_answers:
        return
    # Build lookup: question -> {answer -> value}
    q_data: dict[str, dict[str, float]] = {}
    for q_text, info in group_questions:
        df = info["df"].copy()
        df = df[~df["_answer_"].apply(_is_utility_row)].copy()
        answer_vals: dict[str, float] = {}
        for _, row in df.iterrows():
            ans = str(row["_answer_"]).strip()
            if first_col in df.columns:
                val = float(row[first_col]) if pd.notna(row[first_col]) else 0.0
            else:
                val = 0.0
            answer_vals[ans] = val
        q_data[q_text] = answer_vals
    # Categories = question names (reversed for correct visual order)
    question_names = [q for q, _ in group_questions]
    display_questions = list(reversed(question_names))
    chart_data = CategoryChartData()
    chart_data.categories = display_questions
    for ans in all_answers:
        vals = [q_data.get(q, {}).get(ans, 0.0) for q in display_questions]
        chart_data.add_series(ans, vals)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100,
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(6.7),
        chart_data,
    )
    chart = chart_frame.chart
    n_values = list(dict.fromkeys(info["n_value"] for _, info in group_questions))
    n_display = n_values[0] if len(n_values) == 1 else "/".join(n_values)
    _set_chart_title(chart, f"Groep: {group_id}", n_display, group_id)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.overlap = 100
    _clean_axes(chart)
    if chart.category_axis:
        cat_font = chart.category_axis.tick_labels.font
        _set_font(cat_font, size=Pt(10), color=DARK_GREY)
    num_answers = len(all_answers)
    for s_idx, series in enumerate(plot.series):
        label = all_answers[s_idx] if s_idx < num_answers else ""
        matched = _match_stacked_color(label.lower().strip())
        color = matched if matched else _get_stacked_position_color(s_idx, num_answers)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = hex_to_rgb(color)
        series.has_data_labels = True
        dl = series.data_labels
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False
        dl.font.name = FONT_NAME
        dl.font.size = Pt(10)
        dl.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        dl.number_format = '0"%"'
        dl.number_format_is_linked = False
        dl.position = XL_LABEL_POSITION.CENTER
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _set_font(chart.legend.font, size=Pt(10))
    if significantie_tekst:
        _add_significance_textbox(slide, significantie_tekst)
def _build_grouped_bar_slide(prs, group_questions: list[tuple[str, dict]],
                              selected_cols: list[str], group_id: str,
                              grey_labels: set[str] | None = None,
                              significantie_tekst: str | None = None):
    """Grouped bar chart: shared answer options as categories,
    each question becomes a separate series.  Uses first selected column."""
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    first_col = selected_cols[0] if selected_cols else None
    if not first_col:
        return
    # Collect all unique answer options
    all_answers: list[str] = []
    seen_answers: set[str] = set()
    for _, info in group_questions:
        df = info["df"].copy()
        df = df[~df["_answer_"].apply(_is_utility_row)].copy()
        bottom_set = grey_labels if grey_labels else BOTTOM_LABELS
        df = _sort_bar_df(df, [first_col], grey_labels=bottom_set)
        for ans in df["_answer_"].str.strip().tolist():
            if ans and ans not in seen_answers:
                all_answers.append(ans)
                seen_answers.add(ans)
    if not all_answers:
        return
    chart_data = CategoryChartData()
    chart_data.categories = all_answers
    # Each question becomes a series (reversed for correct visual order)
    series_questions = list(reversed(group_questions))
    for q_text, info in series_questions:
        df = info["df"].copy()
        df = df[~df["_answer_"].apply(_is_utility_row)].copy()
        answer_vals: dict[str, float] = {}
        for _, row in df.iterrows():
            ans = str(row["_answer_"]).strip()
            if first_col in df.columns:
                val = float(row[first_col]) if pd.notna(row[first_col]) else 0.0
            else:
                val = 0.0
            answer_vals[ans] = val
        vals = [answer_vals.get(ans, 0.0) for ans in all_answers]
        chart_data.add_series(q_text, vals)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(6.7),
        chart_data,
    )
    chart = chart_frame.chart
    n_values = list(dict.fromkeys(info["n_value"] for _, info in group_questions))
    n_display = n_values[0] if len(n_values) == 1 else "/".join(n_values)
    _set_chart_title(chart, f"Groep: {group_id}", n_display, group_id)
    plot = chart.plots[0]
    plot.gap_width = 60
    _clean_axes(chart)
    if chart.category_axis:
        cat_font = chart.category_axis.tick_labels.font
        _set_font(cat_font, size=Pt(10), color=DARK_GREY)
    colors = [BAR_PRIMARY, BAR_SECONDARY, "#FF6B81", "#A855F7"]
    num_series = len(series_questions)
    bottom_set = grey_labels if grey_labels else BOTTOM_LABELS
    answers_lower = [a.strip().lower() for a in all_answers]
    for s_idx, series in enumerate(plot.series):
        color_idx = num_series - 1 - s_idx
        base_color = colors[color_idx % len(colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = hex_to_rgb(base_color)
        for pt_idx, ans in enumerate(answers_lower):
            if ans in bottom_set:
                pt = series.points[pt_idx]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = hex_to_rgb(BAR_GREY)
        series.has_data_labels = True
        dl = series.data_labels
        dl.show_value = True
        dl.show_category_name = False
        dl.show_series_name = False
        dl.font.name = FONT_NAME
        dl.font.size = Pt(10)
        dl.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        dl.number_format = '0"%"'
        dl.number_format_is_linked = False
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _set_font(chart.legend.font, size=Pt(10))
    if significantie_tekst:
        _add_significance_textbox(slide, significantie_tekst)
def generate_pptx(questions_data: OrderedDict, config_df: pd.DataFrame,
                  selected_cols: list[str], template_file=None,
                  significantie_aan: bool = False) -> bytes:
    if template_file is not None:
        prs = Presentation(template_file)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
    # If significantie is enabled globally, generate texts now (with Gemini)
    if significantie_aan:
        for q_key, info in questions_data.items():
            if info.get("heeft_significantie") and info.get("significantie_bullets"):
                if "significantie_tekst" not in info or not info["significantie_tekst"]:
                    sig_text = generate_significance_text(q_key, info["significantie_bullets"])
                    if sig_text is None:
                        sig_text = " ".join(info["significantie_bullets"])
                    info["significantie_tekst"] = sig_text
    processed_groups: set[str] = set()
    for _, row in config_df.iterrows():
        if not row.get("Exporteren", False):
            continue
        q_text = row["Vraag"]
        chart_type = row["Grafiektype"]
        group_id = str(row.get("Groep_ID", "")).strip()
        grey_label = str(row.get("Grijs (onderaan)", "")).strip()
        if q_text not in questions_data:
            continue
        info = questions_data[q_text]
        # Determine significance text for this slide
        sig_tekst = None
        if significantie_aan and info.get("heeft_significantie"):
            sig_tekst = info.get("significantie_tekst")
        if group_id:
            # Grouped slide: only generate once per group_id
            if group_id in processed_groups:
                continue
            processed_groups.add(group_id)
            # Collect all exported questions with this group_id + their grey labels
            group_questions: list[tuple[str, dict]] = []
            group_grey: set[str] = set()
            group_sig_parts: list[str] = []
            for _, r2 in config_df.iterrows():
                if not r2.get("Exporteren", False):
                    continue
                g2 = str(r2.get("Groep_ID", "")).strip()
                q2 = r2["Vraag"]
                if g2 == group_id and q2 in questions_data:
                    group_questions.append((q2, questions_data[q2]))
                    gl2 = str(r2.get("Grijs (onderaan)", "")).strip()
                    if gl2:
                        group_grey.add(gl2.lower())
                    if significantie_aan:
                        q2_info = questions_data[q2]
                        if q2_info.get("heeft_significantie") and q2_info.get("significantie_tekst"):
                            group_sig_parts.append(q2_info["significantie_tekst"])
            group_sig_tekst = " ".join(group_sig_parts) if group_sig_parts else None
            if len(group_questions) <= 1:
                # Single question in group -> treat as individual
                if chart_type == "100% Gestapeld horizontaal":
                    _build_stacked_slide(prs, q_text, info, selected_cols, group_id,
                                         significantie_tekst=sig_tekst)
                else:
                    _build_bar_slide(prs, q_text, info, selected_cols, group_id,
                                     grey_label=grey_label, significantie_tekst=sig_tekst)
            else:
                if chart_type == "100% Gestapeld horizontaal":
                    _build_grouped_stacked_slide(prs, group_questions, selected_cols, group_id,
                                                  significantie_tekst=group_sig_tekst)
                else:
                    _build_grouped_bar_slide(prs, group_questions, selected_cols, group_id,
                                             grey_labels=group_grey if group_grey else None,
                                             significantie_tekst=group_sig_tekst)
        else:
            if chart_type == "100% Gestapeld horizontaal":
                _build_stacked_slide(prs, q_text, info, selected_cols, group_id,
                                     significantie_tekst=sig_tekst)
            else:
                _build_bar_slide(prs, q_text, info, selected_cols, group_id,
                                 grey_label=grey_label, significantie_tekst=sig_tekst)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
# =====================================================================
# 4. STREAMLIT APP
# =====================================================================
def _ruigrok_css() -> str:
    """Return all custom CSS for the Ruigrok-branded Streamlit app."""
    return """
    <style>
        /* ── Hide default Streamlit header ── */
        header[data-testid="stHeader"] { background-color: transparent; }
        /* ── Branded top bar ── */
        .ruigrok-bar {
            background: linear-gradient(135deg, #0D0D4F 0%, #1a1a6e 50%, #C60651 100%);
            padding: 1.1rem 2rem;
            display: flex;
            align-items: center;
            gap: 1.2rem;
            border-radius: 0 0 12px 12px;
            margin: -1rem -1rem 1.5rem -1rem;
            box-shadow: 0 4px 16px rgba(13,13,79,0.25);
        }
        .ruigrok-bar svg { height: 40px; width: auto; flex-shrink: 0; }
        .ruigrok-bar .rb-title {
            color: #fff; font-size: 1.55rem; font-weight: 700; letter-spacing: .5px;
        }
        .ruigrok-bar .rb-sub {
            color: rgba(255,255,255,.72); font-size: .88rem; font-weight: 400;
        }
        /* ── Buttons (primary = magenta) ── */
        .stButton > button,
        [data-testid="stBaseButton-secondary"] {
            background-color: #C60651 !important;
            color: white !important;
            border: none !important;
            border-radius: 8px;
            font-weight: 600;
            padding: .5rem 1.4rem;
            transition: all .2s ease;
            box-shadow: 0 2px 8px rgba(198,6,81,.3);
        }
        .stButton > button:hover,
        [data-testid="stBaseButton-secondary"]:hover {
            background-color: #0D0D4F !important;
            color: white !important;
            box-shadow: 0 2px 12px rgba(13,13,79,.4);
            transform: translateY(-1px);
        }
        /* Download button = green */
        .stDownloadButton > button {
            background-color: #39B54A !important;
            color: white !important;
            border: none !important;
            border-radius: 8px;
            font-weight: 600;
            box-shadow: 0 2px 8px rgba(57,181,74,.3);
        }
        .stDownloadButton > button:hover {
            background-color: #2d9a3e !important;
            color: white !important;
        }
        /* ── File uploader "Browse files" button ── */
        [data-testid="stFileUploader"] button {
            background-color: #0D0D4F !important;
            color: white !important;
            border: none !important;
            border-radius: 8px;
            font-weight: 600;
        }
        [data-testid="stFileUploader"] button:hover {
            background-color: #C60651 !important;
            color: white !important;
        }
        /* ── Sidebar ── */
        [data-testid="stSidebar"] {
            background: linear-gradient(180deg, #F8F8FC 0%, #F0F0F8 100%);
            border-right: 2px solid #E8E8F0;
        }
        [data-testid="stSidebar"] h1,
        [data-testid="stSidebar"] h2,
        [data-testid="stSidebar"] h3 { color: #0D0D4F; }
        /* ── General typography ── */
        h1, h2, h3 { color: #0D0D4F; }
        /* ── Cards ── */
        .section-card {
            background: white; border-radius: 10px; padding: 1.2rem;
            border: 1px solid #E8E8F0;
            box-shadow: 0 1px 4px rgba(0,0,0,.06); margin-bottom: 1rem;
        }
        /* ── Data editor ── */
        [data-testid="stDataEditor"] {
            border-radius: 10px; overflow: hidden;
            box-shadow: 0 1px 6px rgba(0,0,0,.08);
        }
        .stAlert { border-radius: 8px; }
        hr { border-color: #E8E8F0; }
        /* ── Stat boxes ── */
        .stat-box {
            background: linear-gradient(135deg, #f8f0f4 0%, #fff 100%);
            border-left: 4px solid #C60651; border-radius: 8px;
            padding: .8rem 1rem; margin-bottom: .5rem;
        }
        .stat-box .label { color: #666; font-size: .8rem; }
        .stat-box .value { color: #0D0D4F; font-size: 1.4rem; font-weight: 700; }
        /* ── Multiselect / selectbox chips ── */
        span[data-baseweb="tag"] {
            background-color: #0D0D4F !important;
        }
    </style>
    """
def _get_logo_html() -> str:
    """Return logo HTML: use logo.png if present, otherwise inline SVG fallback."""
    logo_path = Path(__file__).parent / "logo.png"
    if logo_path.exists():
        b64 = base64.b64encode(logo_path.read_bytes()).decode()
        return f'<img src="data:image/png;base64,{b64}" style="height:42px;" alt="Ruigrok">'
    # Fallback: inline SVG
    return (
        '<svg viewBox="0 0 200 48" xmlns="http://www.w3.org/2000/svg" style="height:42px;">'
        '<circle cx="24" cy="24" r="22" fill="none" stroke="#C60651" stroke-width="3"/>'
        '<path d="M24 6 A18 18 0 0 1 42 24" fill="none" stroke="#fff" stroke-width="3" stroke-linecap="round"/>'
        '<path d="M24 12 A12 12 0 0 1 36 24" fill="none" stroke="rgba(255,255,255,.6)" stroke-width="2.5" stroke-linecap="round"/>'
        '<path d="M24 18 A6 6 0 0 1 30 24" fill="none" stroke="rgba(255,255,255,.35)" stroke-width="2" stroke-linecap="round"/>'
        '<text x="54" y="22" fill="#fff" font-family="Arial,sans-serif" font-size="20" font-weight="700">Ruigrok</text>'
        '<text x="54" y="38" fill="rgba(255,255,255,.7)" font-family="Arial,sans-serif" font-size="11">onderzoek &amp; advies</text>'
        '</svg>'
    )
def main():
    st.set_page_config(page_title="Ruigrok - Grafiek Builder", layout="wide")
    # ── Inject CSS ──
    st.markdown(_ruigrok_css(), unsafe_allow_html=True)
    # ── Branded top bar ──
    st.markdown(
        f'<div class="ruigrok-bar">'
        f'  {_get_logo_html()}'
        f'  <div>'
        f'    <div class="rb-title">Grafiek Builder</div>'
        f'    <div class="rb-sub">Excel &rarr; PowerPoint rapportage</div>'
        f'  </div>'
        f'</div>',
        unsafe_allow_html=True,
    )
    with st.sidebar:
        st.header("1. Upload Excel")
        uploaded = st.file_uploader("Kies een .xlsx bestand", type=["xlsx"])
        st.header("Template (optioneel)")
        template = st.file_uploader(
            "Upload een .pptx template",
            type=["pptx"],
            help="Optioneel: upload een PowerPoint-template met jouw huisstijl en hulplijnen.",
        )
        if template:
            st.session_state.template_file = io.BytesIO(template.read())
        elif "template_file" not in st.session_state:
            st.session_state.template_file = None
        if uploaded:
            if ("uploaded_name" not in st.session_state
                    or st.session_state.uploaded_name != uploaded.name):
                with st.spinner("Excel wordt ingelezen..."):
                    try:
                        questions, data_cols = parse_excel(uploaded)
                    except Exception as e:
                        st.error(f"Fout bij het inlezen: {e}")
                        import traceback
                        st.code(traceback.format_exc())
                        return
                st.session_state.questions = questions
                st.session_state.uploaded_name = uploaded.name
                st.session_state.available_cols = data_cols
                # Check if any question has significance data
                has_any_sig = any(
                    info.get("heeft_significantie", False)
                    for info in questions.values()
                )
                st.session_state.has_any_significantie = has_any_sig
                config_rows = []
                all_answer_options_set: set[str] = set()
                for q_text, info in questions.items():
                    ctype = detect_chart_type(info["answer_options"])
                    grey_auto = ""
                    for a in info["answer_options"]:
                        if a.strip().lower() in BOTTOM_LABELS:
                            grey_auto = a.strip()
                            break
                    config_rows.append({
                        "Exporteren": True,
                        "Vraag": q_text,
                        "Grafiektype": ctype,
                        "Basis (n=)": f"n={info['n_value']}",
                        "Groep_ID": "",
                        "Grijs (onderaan)": grey_auto,
                    })
                    all_answer_options_set.update(
                        a.strip() for a in info["answer_options"] if a.strip()
                    )
                st.session_state.all_answer_options = sorted(all_answer_options_set)
                st.session_state.config_df = pd.DataFrame(config_rows)
                st.success(f"{len(questions)} vragen gevonden!")
        if "available_cols" in st.session_state and st.session_state.available_cols:
            st.header("2. Kolommen")
            selected = st.multiselect(
                "Selecteer uitsplitsingen",
                options=st.session_state.available_cols,
                default=st.session_state.available_cols[:1],
                help="Kies welke kolommen (Totaal, Man, Vrouw, ...) je in de grafieken wilt.",
            )
            st.session_state.selected_cols = selected
        # Global significantie toggle in sidebar
        if st.session_state.get("has_any_significantie", False):
            st.header("3. Significantie")
            sig_toggle = st.toggle(
                "Significantie tonen",
                value=st.session_state.get("significantie_aan", False),
                key="sig_toggle",
                help="Toon significante verschillen onder alle grafieken waarvoor data beschikbaar is.",
            )
            st.session_state.significantie_aan = sig_toggle
            if sig_toggle:
                sig_count = sum(
                    1 for info in st.session_state.get("questions", {}).values()
                    if info.get("heeft_significantie", False)
                )
                st.caption(f"Significantiedata gevonden voor {sig_count} vragen.")
    if "config_df" not in st.session_state:
        st.markdown(
            '<div class="section-card" style="text-align:center;padding:3rem;">'
            '<h3 style="margin-bottom:.5rem;">Welkom bij de Grafiek Builder</h3>'
            '<p style="color:#666;">Upload een Excel-bestand via de sidebar om te beginnen '
            'met het genereren van PowerPoint rapportages.</p>'
            '</div>',
            unsafe_allow_html=True,
        )
        return
    st.header("Regie Tabel")
    st.caption("Pas instellingen per vraag aan. Vink 'Exporteren' aan voor opname in de PowerPoint.")
    # Show main config table WITHOUT the grey column (handled separately below)
    display_df = st.session_state.config_df.drop(columns=["Grijs (onderaan)"], errors="ignore")
    edited = st.data_editor(
        display_df,
        column_config={
            "Exporteren": st.column_config.CheckboxColumn("Exporteren", default=True),
            "Vraag": st.column_config.TextColumn("Vraag", disabled=True, width="large"),
            "Grafiektype": st.column_config.SelectboxColumn(
                "Grafiektype",
                options=["Staafdiagram", "100% Gestapeld horizontaal"],
                width="medium",
            ),
            "Basis (n=)": st.column_config.TextColumn("Basis (n=)", disabled=True, width="small"),
            "Groep_ID": st.column_config.TextColumn("Groep_ID", width="small"),
        },
        use_container_width=True,
        num_rows="fixed",
        hide_index=True,
        key="regie_editor",
    )
    # ── Per-vraag "Grijs (onderaan)" selectie ──
    # Only show for exported questions; each gets its own answer options
    questions = st.session_state.questions
    grey_values = st.session_state.config_df["Grijs (onderaan)"].tolist() if "Grijs (onderaan)" in st.session_state.config_df.columns else [""] * len(edited)
    with st.expander("Alternatieve 'weet ik niet' selecteren (grijze kleur)", expanded=False):
        st.caption("Kies per vraag welk antwoord grijs en onderaan de grafiek moet worden weergegeven.")
        new_grey = []
        for idx, row in edited.iterrows():
            q_text = row["Vraag"]
            if not row.get("Exporteren", False):
                new_grey.append(grey_values[idx] if idx < len(grey_values) else "")
                continue
            # Get answer options for THIS specific question
            q_options = ["(geen)"]
            if q_text in questions:
                q_options += questions[q_text]["answer_options"]
            current_grey = grey_values[idx] if idx < len(grey_values) else ""
            # Find index of current value in options
            default_idx = 0
            if current_grey and current_grey in q_options:
                default_idx = q_options.index(current_grey)
            selected_grey = st.selectbox(
                q_text,
                options=q_options,
                index=default_idx,
                key=f"grey_{idx}",
                label_visibility="visible",
            )
            new_grey.append("" if selected_grey == "(geen)" else selected_grey)
    # Merge grey values back into config_df
    merged = edited.copy()
    merged["Grijs (onderaan)"] = new_grey
    st.session_state.config_df = merged
    # ── Debug info (collapsed, alleen zichtbaar in debug mode) ──
    if DEBUG_MODE:
        with st.expander("Debug: data per vraag"):
            questions = st.session_state.questions
            q_choice = st.selectbox("Kies een vraag", list(questions.keys()), key="debug_q")
            if q_choice:
                info = questions[q_choice]
                st.write(f"**n = {info['n_value']}**")
                st.write(f"**Antwoordopties ({len(info['answer_options'])}):** {info['answer_options']}")
                display_df = info["df"].drop(columns=["_question_"], errors="ignore")
                st.dataframe(display_df, use_container_width=True, hide_index=True)
    st.divider()
    selected_cols = st.session_state.get("selected_cols", [])
    if not selected_cols:
        st.warning("Selecteer minstens een kolom in de sidebar.")
        return
    export_count = int(edited["Exporteren"].sum())
    # Stats overview
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown(f'<div class="stat-box"><div class="label">Totaal vragen</div><div class="value">{len(edited)}</div></div>', unsafe_allow_html=True)
    with col2:
        st.markdown(f'<div class="stat-box"><div class="label">Te exporteren</div><div class="value">{export_count}</div></div>', unsafe_allow_html=True)
    with col3:
        st.markdown(f'<div class="stat-box"><div class="label">Kolommen</div><div class="value">{len(selected_cols)}</div></div>', unsafe_allow_html=True)
    st.write("")
    if st.button(
        f"Genereer PowerPoint ({export_count} slides)",
        type="primary",
        disabled=(export_count == 0),
    ):
        sig_aan = st.session_state.get("significantie_aan", False)
        with st.spinner("PowerPoint wordt gegenereerd..."):
            try:
                tpl = st.session_state.get("template_file")
                if tpl:
                    tpl.seek(0)
                pptx_bytes = generate_pptx(
                    st.session_state.questions, st.session_state.config_df, selected_cols,
                    template_file=tpl,
                    significantie_aan=sig_aan,
                )
            except Exception as e:
                st.error(f"Fout bij genereren: {e}")
                import traceback
                st.code(traceback.format_exc())
                return
        st.success("PowerPoint is klaar!")
        st.download_button(
            label="Download rapportage.pptx",
            data=pptx_bytes,
            file_name="rapportage.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
if __name__ == "__main__":
    main()
